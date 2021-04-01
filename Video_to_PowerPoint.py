#!/usr/bin/env python
# coding: utf-8

# In[1]:


# Take screeshot of designated coordinates
get_ipython().system('pip install pyscreenshot')


# In[2]:


# Identifies coordinates to take screenshot
get_ipython().system('pip install pynput')


# In[3]:


# Creates PPTX file and saves image to slide
get_ipython().system('pip install python-pptx')


# In[4]:


# Identifies changes between two images
get_ipython().system('pip install scikit-image')


# In[5]:


# Use pynput.mouse.Listener
from pynput.mouse import Listener
from pynput import keyboard
import time

#Grab screenshot of selected screen coordinates
import pyscreenshot as ImageGrab

# Inserting image into the PowerPoint file
from pptx import Presentation 
from pptx.util import Inches

import numpy as np
import matplotlib.pyplot as plt

# Calculate mean squared error between images
from skimage import data, img_as_float
from skimage.metrics import mean_squared_error


# In[9]:


# Set name of powerpoint file
pptx_file = "Insert file name here"


# In[10]:


# Capture screen area of mouse click and release
# Must click in upper left corner and release in lower right corner of area to capture

down_x = down_y = up_x = up_y = -1

def on_click(x, y, button, pressed):
    global down_x
    global down_y
    global up_x
    global up_y
    if pressed:
        (down_x, down_y) = (x, y)
    else:
        (up_x, up_y) = (x, y)
        return False

with Listener(on_click=on_click) as listener:
    listener.join()

print("Mouse drag from", down_x, ",", down_y, "to", up_x, ",", up_y)


# In[12]:


# Press 'End' key to close process

break_program = False
def on_press(key):
    global break_program
    #print (key)
    if key == keyboard.Key.end:
        print ('End Pressed')
        break_program = True
        return False


# part of the screen
im = ImageGrab.grab(bbox=(down_x, down_y, up_x, up_y))  # X1,Y1,X2,Y2
# save im as im1
im1 = im
# save first screenshot file
im.save("Images/image_1.png")
# Giving Image path 
img_path = "Images/image_1.png"
# Creating an Presentation object
ppt = Presentation() 
# Selecting blank slide
blank_slide_layout = ppt.slide_layouts[6]
# Attaching slide to ppt
slide = ppt.slides.add_slide(blank_slide_layout) 
# For margins
left = top = Inches(1) 
left = Inches(.5) 
height = Inches(5) 
pic = slide.shapes.add_picture(img_path, left, top, height = height)
# save file
ppt.save('PowerPoint/'+ pptx_file + '.pptx')

    
with keyboard.Listener(on_press=on_press) as listener:
    while break_program == False:
        im2 = ImageGrab.grab(bbox=(down_x, down_y, up_x, up_y))  # X1,Y1,X2,Y2
        im2.save("Images/image_2.png")
        # Compare images using Mean Squared Error
        img1 = img_as_float(im1)
        img2 = img_as_float(im2)
        mse = mean_squared_error(img1, img2)
        # Giving Image path 
        img_path = "Images/image_2.png"
        # Creating an Presentation object
        ppt = Presentation('PowerPoint/'+ pptx_file + '.pptx') 
        # Selecting blank slide
        blank_slide_layout = ppt.slide_layouts[6] 
        # Attaching slide to ppt
        slide = ppt.slides.add_slide(blank_slide_layout) 
        # For margins
        left = top = Inches(1) 
        left = Inches(.5) 
        height = Inches(5) 
        pic = slide.shapes.add_picture(img_path, left, top, height = height)
        # if mean squared error is greater that .005, append screenshot to pptx file
        if mse > 0.005:
            ppt.save('PowerPoint/'+ pptx_file + '.pptx')
            im1 = im2
        time.sleep(5)
    listener.join()


# In[ ]:




